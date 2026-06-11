from __future__ import annotations

import re
import shutil
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SUBMISSION_DIR = ROOT / "submission_docs"
SOURCE_PPT = SUBMISSION_DIR / "generative agents.pptx"
OUTPUT_PPT = SUBMISSION_DIR / "generative agents_最终版.pptx"

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

ET.register_namespace("p", P_NS)
ET.register_namespace("a", A_NS)
ET.register_namespace("r", R_NS)

SLIDE_PAYLOADS = [
    {
        "title": "复现机制总览",
        "subtitle": "论文复现补充",
        "section": "核心闭环",
        "cards": [
            "观察\n读取环境与相遇",
            "记忆\n写入经历",
            "检索\n召回相关记忆",
            "反思\n形成高层认知",
        ],
        "detail": "环境感知 -> 记忆检索 -> 计划驱动 -> 行动/对话 -> 记忆更新 -> 反思形成\n本项目保留论文中最核心的认知闭环，角色行为由计划、位置、记忆和相遇共同驱动。",
    },
    {
        "title": "双模式设计",
        "subtitle": "论文复现补充",
        "section": "运行模式",
        "cards": [
            "deterministic\n稳定演示",
            "llm\n模型生成",
            "fallback\n失败回退",
            "evidence\n状态可解释",
        ],
        "detail": "deterministic 模式适合录屏和答辩；llm 模式让模型参与计划、行动、对话和 reflection 生成。\n接口失败时自动规则回退，右侧面板显示模式、LLM 状态、记忆流数量和反思触发依据。",
    },
    {
        "title": "LLM 模式演示路线",
        "subtitle": "论文复现补充",
        "section": "时间线",
        "cards": [
            "08:00\n初始态",
            "10:00\nAlice -> Bob",
            "14:00\nBob -> Carol",
            "14:30\nReflection",
        ],
        "detail": "08:00 展示 persona、初始记忆和计划；10:00 展示对话如何改写 Bob 内部状态；14:00 展示局部互动形成传播链；14:30 展示记忆重要性超过阈值后的高层反思。",
    },
    {
        "title": "对应论文机制",
        "subtitle": "论文复现补充",
        "section": "机制映射",
        "cards": [
            "Memory Stream\n统一记忆流",
            "Retrieval\n检索函数",
            "Planning\n计划行动",
            "Reflection\n抽象总结",
        ],
        "detail": "观察、行动、对话和反思都写入 memory stream；检索按重要性、近因、地点、社交关系和相关性打分；计划提供日程骨架，检索记忆影响当前行动。",
    },
    {
        "title": "局限与改进",
        "subtitle": "论文复现补充",
        "section": "后续方向",
        "cards": [
            "规模\n3 个角色",
            "世界\n4 个地点",
            "检索\n待接 embedding",
            "计划\n待扩展多日",
        ],
        "detail": "当前是课程作业级复现，不是 Stanford Smallville 全量工程。\n后续可以扩展 25 个 agent、embedding retrieval、多日长期计划、更丰富社交关系和开放式事件解析。",
    },
]


def qname(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


def slide_number(name: str) -> int:
    match = re.search(r"slides/slide(\d+)\.xml$", name)
    return int(match.group(1)) if match else -1


def next_rel_id(root: ET.Element) -> str:
    max_id = 0
    for rel in root:
        rid = rel.attrib.get("Id", "")
        if rid.startswith("rId") and rid[3:].isdigit():
            max_id = max(max_id, int(rid[3:]))
    return f"rId{max_id + 1}"


def max_slide_id(presentation: ET.Element) -> int:
    values = []
    for item in presentation.findall(f".//{{{P_NS}}}sldId"):
        value = item.attrib.get("id")
        if value and value.isdigit():
            values.append(int(value))
    return max(values or [255])


def replace_slide_text(slide_xml: bytes, lines: list[str]) -> bytes:
    root = ET.fromstring(slide_xml)
    text_nodes = root.findall(f".//{{{A_NS}}}t")
    if not text_nodes:
        return slide_xml

    for index, node in enumerate(text_nodes):
        node.text = lines[index] if index < len(lines) else ""

    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def placeholder_lines(payload: dict[str, object]) -> list[str]:
    cards = payload["cards"]
    assert isinstance(cards, list)
    return [
        str(payload["title"]),
        str(payload["subtitle"]),
        str(payload["section"]),
        *[str(item) for item in cards],
        str(payload["detail"]),
    ]


def add_content_type(content_types_xml: bytes, slide_path: str) -> bytes:
    root = ET.fromstring(content_types_xml)
    part_name = f"/{slide_path}"
    exists = any(item.attrib.get("PartName") == part_name for item in root)
    if not exists:
        ET.SubElement(
            root,
            qname(CT_NS, "Override"),
            {
                "PartName": part_name,
                "ContentType": "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
            },
        )
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def extend_presentation_xml(presentation_xml: bytes, new_slides: list[tuple[int, str]]) -> bytes:
    root = ET.fromstring(presentation_xml)
    slide_id_list = root.find(f"{{{P_NS}}}sldIdLst")
    if slide_id_list is None:
        raise ValueError("presentation.xml does not contain p:sldIdLst")

    current_max_id = max_slide_id(root)
    for index, (_slide_no, rid) in enumerate(new_slides, start=1):
        ET.SubElement(
            slide_id_list,
            qname(P_NS, "sldId"),
            {"id": str(current_max_id + index), qname(R_NS, "id"): rid},
        )
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def extend_presentation_rels(rels_xml: bytes, new_slides: list[tuple[int, str]]) -> bytes:
    root = ET.fromstring(rels_xml)
    for slide_no, rid in new_slides:
        ET.SubElement(
            root,
            qname(REL_NS, "Relationship"),
            {
                "Id": rid,
                "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
                "Target": f"slides/slide{slide_no}.xml",
            },
        )
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def build_final_ppt() -> None:
    if not SOURCE_PPT.exists():
        raise FileNotFoundError(f"Missing group PPT: {SOURCE_PPT}")

    OUTPUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(SOURCE_PPT, OUTPUT_PPT)

    with zipfile.ZipFile(OUTPUT_PPT, "r") as src:
        entries = {name: src.read(name) for name in src.namelist()}

    slide_names = sorted(
        [name for name in entries if re.search(r"ppt/slides/slide\d+\.xml$", name)],
        key=slide_number,
    )
    if not slide_names:
        raise ValueError("No slides found in group PPT")

    base_slide_name = slide_names[-1]
    base_slide_no = slide_number(base_slide_name)
    base_rels_name = f"ppt/slides/_rels/slide{base_slide_no}.xml.rels"
    base_slide_xml = entries[base_slide_name]
    base_rels_xml = entries.get(base_rels_name)

    presentation_rels_root = ET.fromstring(entries["ppt/_rels/presentation.xml.rels"])
    new_slide_refs: list[tuple[int, str]] = []

    for offset, payload in enumerate(SLIDE_PAYLOADS, start=1):
        new_slide_no = base_slide_no + offset
        new_slide_name = f"ppt/slides/slide{new_slide_no}.xml"
        new_rels_name = f"ppt/slides/_rels/slide{new_slide_no}.xml.rels"
        rid = next_rel_id(presentation_rels_root)
        ET.SubElement(
            presentation_rels_root,
            qname(REL_NS, "Relationship"),
            {
                "Id": rid,
                "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
                "Target": f"slides/slide{new_slide_no}.xml",
            },
        )
        new_slide_refs.append((new_slide_no, rid))
        entries[new_slide_name] = replace_slide_text(base_slide_xml, placeholder_lines(payload))
        if base_rels_xml:
            entries[new_rels_name] = base_rels_xml
        entries["[Content_Types].xml"] = add_content_type(entries["[Content_Types].xml"], new_slide_name)

    # presentation_rels_root has already been extended so serialize it directly.
    entries["ppt/_rels/presentation.xml.rels"] = ET.tostring(
        presentation_rels_root,
        encoding="utf-8",
        xml_declaration=True,
    )
    entries["ppt/presentation.xml"] = extend_presentation_xml(entries["ppt/presentation.xml"], new_slide_refs)

    with zipfile.ZipFile(OUTPUT_PPT, "w", zipfile.ZIP_DEFLATED) as dst:
        for name, data in entries.items():
            dst.writestr(name, data)

    polish_added_slides()
    print(f"Wrote {OUTPUT_PPT.relative_to(ROOT)}")


def set_shape_text(slide, index: int, text: str) -> None:
    if index >= len(slide.shapes):
        return
    shape = slide.shapes[index]
    if hasattr(shape, "text"):
        shape.text = text


def polish_added_slides() -> None:
    presentation = Presentation(str(OUTPUT_PPT))
    start = len(presentation.slides) - len(SLIDE_PAYLOADS)
    shape_map = {
        "title": 2,
        "subtitle": 3,
        "section": 4,
        "cards": [8, 9, 10, 11],
        "detail": 12,
    }
    for offset, payload in enumerate(SLIDE_PAYLOADS):
        slide = presentation.slides[start + offset]
        set_shape_text(slide, shape_map["title"], str(payload["title"]))
        set_shape_text(slide, shape_map["subtitle"], str(payload["subtitle"]))
        set_shape_text(slide, shape_map["section"], str(payload["section"]))
        for index, text in zip(shape_map["cards"], payload["cards"]):
            set_shape_text(slide, index, str(text))
        set_shape_text(slide, shape_map["detail"], str(payload["detail"]))
    presentation.save(str(OUTPUT_PPT))


if __name__ == "__main__":
    build_final_ppt()

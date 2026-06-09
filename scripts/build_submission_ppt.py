from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\demo\buaa\suanfa\three")
TEMPLATE = ROOT / "作业模板" / "作业模板" / "解题ppt.pptx"
OUTPUT = ROOT / "submission_docs" / "Generative_Agents_课程答辩稿.pptx"

ACCENT = RGBColor(199, 76, 44)
DARK = RGBColor(37, 43, 56)
MUTED = RGBColor(95, 102, 119)
LIGHT = RGBColor(246, 241, 234)
GREEN = RGBColor(102, 143, 118)
BLUE = RGBColor(90, 136, 197)


SLIDES = [
    {
        "type": "title",
        "title": "Generative Agents 论文复刻",
        "subtitle": "基于 Smallville 风格场景的最小完整机制演示系统\n《大模型原理与应用》课程作业",
    },
    {
        "type": "content",
        "title": "为什么要研究 Generative Agents",
        "bullets": [
            "传统 NPC 往往依赖固定脚本，缺乏持续记忆和自主行为。",
            "大语言模型使长期上下文驱动的角色行为生成成为可能。",
            "论文提出了开放小镇环境中的生成式智能体框架。",
            "智能体不仅能行动，还能形成关系、传播信息、积累经验。",
        ],
    },
    {
        "type": "content",
        "title": "论文试图解决什么问题",
        "bullets": [
            "如何让智能体在开放环境中表现出持续、一致、可解释的行为。",
            "如何让角色行为依赖过往记忆，而不是即时随机输出。",
            "如何让多个角色之间形成社会性互动和信息传播。",
            "如何在行为链中引入高层反思，影响未来计划。",
        ],
    },
    {
        "type": "diagram_loop",
        "title": "Generative Agents 的核心闭环",
    },
    {
        "type": "content",
        "title": "我们的复刻目标是什么",
        "bullets": [
            "复刻 Smallville 风格小镇环境。",
            "使用 3 个差异化 NPC。",
            "展示时间推进与地点移动。",
            "实现计划、记忆、检索、对话、反思。",
            "形成一条可解释的信息传播链。",
        ],
    },
    {
        "type": "content",
        "title": "为什么选择最小完整机制方案",
        "bullets": [
            "作业提交时间紧，不能追求论文级高复杂度。",
            "课程答辩更需要可解释、可演示、可运行的系统。",
            "全量复现会把精力分散到大量非关键细节上。",
            "最小完整机制更适合形成代码、报告、PPT 和视频的统一交付链。",
        ],
    },
    {
        "type": "diagram_arch",
        "title": "系统总体方案概览",
    },
    {
        "type": "content",
        "title": "前端展示层设计",
        "bullets": [
            "使用 Next.js + React 构建 2D 小镇可视化界面。",
            "显示 NPC 所在位置、当前动作和简短 utterance。",
            "展示当前时间、计划、记忆、反思和事件日志。",
            "支持点击 NPC 查看详细解释面板。",
        ],
    },
    {
        "type": "content",
        "title": "后端仿真编排层设计",
        "bullets": [
            "使用 FastAPI 暴露状态接口和 WebSocket。",
            "simulation.py 负责推进 world tick 和模拟时间。",
            "每个 tick 更新 active plan、地点、事件和交互。",
            "相遇时触发信息传播与 reflection 逻辑。",
        ],
    },
    {
        "type": "content",
        "title": "Agent Cognition 认知层设计",
        "bullets": [
            "负责记忆检索、计划生成、行动生成、对话生成和反思生成。",
            "将行为生成逻辑从主仿真循环中解耦。",
            "支持 LLM 模式与 fallback 模式双通路。",
            "让系统结构更清晰，也更方便写报告和答辩讲解。",
        ],
    },
    {
        "type": "content",
        "title": "核心数据结构",
        "bullets": [
            "Agent：profile、active plan、recent memories、retrieved memories、reflections。",
            "PlanItem：时间段、地点、任务描述。",
            "MemoryEntry：文本、时间、地点、重要性、相关角色、类型。",
            "WorldState：时间、地点、角色列表、事件日志、LLM 状态。",
        ],
    },
    {
        "type": "content",
        "title": "Daily Plan 是如何生成的",
        "bullets": [
            "系统支持 LLM 生成和静态计划兜底两种模式。",
            "LLM 模式输入 agent profile、地点集合与已有 reflection。",
            "输出为一天内多个时间段的粗粒度计划项。",
            "兜底模式保证没有 API key 时仍能稳定演示。",
        ],
    },
    {
        "type": "content",
        "title": "记忆系统如何工作",
        "bullets": [
            "记忆分为 observation、conversation、reflection 三类。",
            "每条记忆记录文本、时间戳、地点、重要性和相关角色。",
            "conversation memory 来自角色对话和信息传播。",
            "reflection memory 来自更高层的社会性总结。",
        ],
    },
    {
        "type": "content",
        "title": "当前情境下如何检索相关记忆",
        "bullets": [
            "系统采用轻量记忆打分策略，而非复杂向量数据库。",
            "综合考虑重要性、时间新鲜度、地点相关性、附近角色和关键词匹配。",
            "每次行动前检索 Top-K memory 作为上下文。",
            "这种方案够用、可讲、可落地，适合课程作业。",
        ],
    },
    {
        "type": "content",
        "title": "NPC 当前动作如何生成",
        "bullets": [
            "输入包含当前时间、地点、active plan、附近角色和 retrieved memories。",
            "输出包含 summary、utterance 和 reasoning note。",
            "reasoning note 直接用于前端解释面板展示。",
            "这让系统不仅能跑，还能解释为什么这样行动。",
        ],
    },
    {
        "type": "diagram_chain",
        "title": "社交传播链如何形成",
    },
    {
        "type": "content",
        "title": "Reflection 是如何形成的",
        "bullets": [
            "当 gathering 信息在多个角色之间传播后，系统生成高层 reflection。",
            "reflection 会重新写回 memory bank，成为未来决策上下文。",
            "这一步把事件从个体知识提升为共享社会知识。",
            "它是论文机制里最关键也最容易被忽略的一层。",
        ],
    },
    {
        "type": "content",
        "title": "为什么采用 LLM + Fallback 双模式",
        "bullets": [
            "仅使用 LLM 会带来 API 成本、网络依赖和演示稳定性问题。",
            "仅使用规则模板又难以体现课程中的大模型特色。",
            "因此系统采用 LLM 优先、Fallback 兜底的双模式架构。",
            "这是一种兼顾生成性与工程稳定性的折中方案。",
        ],
    },
    {
        "type": "placeholder",
        "title": "系统界面展示",
        "label": "待补截图：地图主界面 + 右侧状态面板",
        "caption": "建议放全屏主界面截图，突出 4 个地点、3 个 NPC、时间、事件日志与控制按钮。",
    },
    {
        "type": "placeholder",
        "title": "为什么这个 NPC 此刻这样行动",
        "label": "待补截图：选中 NPC 的解释面板",
        "caption": "建议突出 Active plan、Retrieved memories、Latest utterance 和 Reasoning context。",
    },
    {
        "type": "content",
        "title": "当前系统已实现的完整行为链",
        "bullets": [
            "角色按计划移动。",
            "根据当前环境检索记忆。",
            "生成当前动作和一句话。",
            "角色相遇并进行信息传播。",
            "生成新的 conversation memory。",
            "在共享知识形成后生成 reflection。",
        ],
    },
    {
        "type": "content",
        "title": "项目中的关键技术难点",
        "bullets": [
            "如何在有限时间内控制复刻范围。",
            "如何避免系统沦为纯脚本剧情。",
            "如何让系统状态可解释、可答辩。",
            "如何兼顾 LLM 特征与本地稳定性。",
            "如何让代码、报告、PPT 和视频形成统一交付链。",
        ],
    },
    {
        "type": "content",
        "title": "本项目的三个工程创新点",
        "bullets": [
            "双模式认知架构：LLM 与 fallback 并存。",
            "面向答辩的 reasoning note：直接暴露当前行为依据。",
            "可控社交传播链：让信息传播与 reflection 稳定出现。",
        ],
    },
    {
        "type": "content",
        "title": "当前系统的局限性与改进方向",
        "bullets": [
            "NPC 数量仍然较少，场景规模有限。",
            "记忆检索仍是轻量规则打分。",
            "计划粒度较粗，尚未支持多天连续仿真。",
            "后续可引入更多 NPC、向量检索和更复杂事件链。",
        ],
    },
    {
        "type": "content",
        "title": "总结",
        "bullets": [
            "本项目完成了课程作业场景下的最小完整机制 Generative Agents 复刻。",
            "系统已经具备环境、计划、记忆、行动、对话、传播和反思闭环。",
            "相比只做 UI 或过度追求论文复杂度，本项目更强调可运行、可解释、可演示、可交付。",
            "谢谢老师。",
        ],
    },
]


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)  # pyright: ignore[reportAttributeAccessIssue]
    for slide_id in slide_ids:
        prs.slides._sldIdLst.remove(slide_id)  # pyright: ignore[reportAttributeAccessIssue]


def get_placeholder_by_type(slide, idx: int):
    placeholders = list(slide.placeholders)
    if idx < len(placeholders):
        return placeholders[idx]
    return None


def style_text_frame(text_frame, font_size=20, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Microsoft YaHei"


def add_title_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = get_placeholder_by_type(slide, 1)
    body.text = bullets[0]
    for bullet in bullets[1:]:
        p = body.text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
    style_text_frame(slide.shapes.title.text_frame, font_size=24, bold=True)
    style_text_frame(body.text_frame, font_size=18)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    sub = get_placeholder_by_type(slide, 1)
    if sub is not None:
        sub.text = subtitle
        style_text_frame(sub.text_frame, font_size=18, color=MUTED, align=PP_ALIGN.CENTER)
    style_text_frame(slide.shapes.title.text_frame, font_size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def add_placeholder_slide(prs: Presentation, title: str, label: str, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_text_frame(slide.shapes.title.text_frame, font_size=24, bold=True)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.6), Inches(8.0), Inches(2.9))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = ACCENT
    tf = box.text_frame
    tf.text = label
    style_text_frame(tf, font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    note = slide.shapes.add_textbox(Inches(1.1), Inches(4.8), Inches(8.0), Inches(1.2))
    note.text_frame.text = caption
    style_text_frame(note.text_frame, font_size=16, color=MUTED)


def add_loop_diagram(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_text_frame(slide.shapes.title.text_frame, font_size=24, bold=True)
    labels = ["感知", "记忆检索", "计划", "行动 / 对话", "记忆更新", "反思"]
    positions = [
        (1.0, 2.6, BLUE),
        (2.7, 1.5, GREEN),
        (5.1, 1.5, ACCENT),
        (6.8, 2.6, BLUE),
        (5.1, 3.8, GREEN),
        (2.7, 3.8, ACCENT),
    ]
    for label, (x, y, color) in zip(labels, positions):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.6), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        shape.text_frame.text = label
        style_text_frame(shape.text_frame, font_size=16, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)

    arrows = [
        ("→", 2.15, 1.95),
        ("→", 4.6, 1.95),
        ("↓", 7.1, 3.0),
        ("←", 4.65, 4.2),
        ("←", 2.15, 4.2),
        ("↑", 1.0, 3.0),
    ]
    for char, x, y in arrows:
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.6), Inches(0.5))
        tb.text_frame.text = char
        style_text_frame(tb.text_frame, font_size=24, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def add_architecture_diagram(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_text_frame(slide.shapes.title.text_frame, font_size=24, bold=True)

    boxes = [
        ("Frontend UI\nMap / Panels / Event Log", 0.6, 1.6, 2.2, 1.0, BLUE),
        ("FastAPI + WebSocket", 3.1, 1.6, 2.0, 1.0, ACCENT),
        ("World Engine", 5.5, 1.6, 1.7, 1.0, GREEN),
        ("Planner", 3.0, 3.0, 1.5, 0.8, BLUE),
        ("Memory System", 4.8, 3.0, 1.8, 0.8, GREEN),
        ("Action Generator", 6.9, 3.0, 1.9, 0.8, ACCENT),
        ("LLM / Fallback", 6.6, 4.4, 2.2, 0.9, DARK),
    ]
    for text, x, y, w, h, color in boxes:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        shape.text_frame.text = text
        style_text_frame(shape.text_frame, font_size=16, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)

    for char, x, y in [
        ("→", 2.75, 1.95),
        ("→", 5.1, 1.95),
        ("↓", 3.65, 2.45),
        ("↓", 5.5, 2.45),
        ("↓", 7.55, 2.45),
        ("↘", 4.6, 3.85),
        ("↓", 7.55, 3.85),
    ]:
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.5), Inches(0.4))
        tb.text_frame.text = char
        style_text_frame(tb.text_frame, font_size=20, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def add_chain_diagram(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_text_frame(slide.shapes.title.text_frame, font_size=24, bold=True)
    nodes = [
        ("Alice\n知道 gathering", 0.7, 2.4, BLUE),
        ("Bob\n在 cafe 获知", 3.1, 2.4, ACCENT),
        ("Carol\n在 square 获知", 5.5, 2.4, GREEN),
        ("Shared\nReflection", 7.8, 2.4, DARK),
    ]
    for text, x, y, color in nodes:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.6), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        shape.text_frame.text = text
        style_text_frame(shape.text_frame, font_size=16, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    for char, x in [("→", 2.45), ("→", 4.9), ("→", 7.2)]:
        tb = slide.shapes.add_textbox(Inches(x), Inches(2.75), Inches(0.5), Inches(0.4))
        tb.text_frame.text = char
        style_text_frame(tb.text_frame, font_size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    note = slide.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(8.0), Inches(0.8))
    note.text_frame.text = "当前演示链路：10:00 Alice -> Bob；14:00 Bob -> Carol；随后形成共享反思。"
    style_text_frame(note.text_frame, font_size=17, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


def build() -> None:
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)

    for slide_data in SLIDES:
        slide_type = slide_data["type"]
        if slide_type == "title":
            add_title_slide(prs, slide_data["title"], slide_data["subtitle"])
        elif slide_type == "content":
            add_title_content_slide(prs, slide_data["title"], slide_data["bullets"])
        elif slide_type == "placeholder":
            add_placeholder_slide(prs, slide_data["title"], slide_data["label"], slide_data["caption"])
        elif slide_type == "diagram_loop":
            add_loop_diagram(prs, slide_data["title"])
        elif slide_type == "diagram_arch":
            add_architecture_diagram(prs, slide_data["title"])
        elif slide_type == "diagram_chain":
            add_chain_diagram(prs, slide_data["title"])
        else:
            raise ValueError(f"Unknown slide type: {slide_type}")

    prs.save(str(OUTPUT))
    print(f"saved: {OUTPUT}")


if __name__ == "__main__":
    build()
